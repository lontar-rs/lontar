#!/usr/bin/env python3
"""Generate DOCX and PPTX reference documents from strings.txt."""
import pathlib
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches

ROOT = pathlib.Path(__file__).parent
STRINGS = ROOT / "strings.txt"
REFERENCE_DIR = ROOT.parent / "reference_docs"
REFERENCE_DIR.mkdir(parents=True, exist_ok=True)
OUT_DOCX = REFERENCE_DIR / "multiscript.docx"
OUT_PPTX = REFERENCE_DIR / "multiscript.pptx"


def parse_strings():
    sections = []
    current = None
    for line in STRINGS.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if not line or line.startswith("#"):
            continue
        if line.startswith("[") and line.endswith("]"):
            current = {"section": line.strip("[]"), "entries": []}
            sections.append(current)
        else:
            if current is None:
                continue
            if ":" in line:
                key, val = line.split(":", 1)
                current["entries"].append((key.strip(), val.strip()))
    return sections


def generate_docx(sections):
    doc = DocxDocument()
    doc.add_heading("Multi-Script Reference", level=1)
    for section in sections:
        doc.add_heading(section["section"], level=2)
        for key, text in section["entries"]:
            p = doc.add_paragraph()
            p.add_run(f"{key}: ")
            p.add_run(text)
    doc.save(OUT_DOCX)


def generate_pptx(sections):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    title_slide_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Multi-Script Reference"
    slide.placeholders[1].text = "Generated from strings.txt"

    for section in sections:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = section["section"]
        body = slide.shapes.placeholders[1].text_frame
        for key, text in section["entries"]:
            p = body.add_paragraph()
            p.text = f"{key}: {text}"
    prs.save(OUT_PPTX)


def main():
    sections = parse_strings()
    generate_docx(sections)
    generate_pptx(sections)
    print(f"Wrote {OUT_DOCX} and {OUT_PPTX}")


if __name__ == "__main__":
    main()
