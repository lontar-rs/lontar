#!/usr/bin/env python3
"""Generate multi-script reference DOCX and PPTX documents.

Outputs:
- tests/fixtures/reference_docs/multiscript.docx
- tests/fixtures/reference_docs/multiscript.pptx

Uses `test_strings.json` for representative text across scripts. Set fonts to
Noto families where appropriate (font files are not embedded).
"""
from __future__ import annotations

import json
from pathlib import Path
from typing import Dict, Iterable, Tuple

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
try:
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ModuleNotFoundError:
    HAS_PPTX = False

HERE = Path(__file__).resolve().parent
FIXTURES = HERE.parent
REFERENCE_DOCS = FIXTURES / "reference_docs"
TEST_STRINGS_PATH = HERE / "test_strings.json"

# Maps ISO 15924 script codes (or combined markers) to preferred Noto font names.
FONT_MAP: Dict[str, str] = {
    "Latn": "Noto Sans",
    "Cyrl": "Noto Sans",
    "Grek": "Noto Sans",
    "Geor": "Noto Sans Georgian",
    "Deva": "Noto Sans Devanagari",
    "Beng": "Noto Sans Bengali",
    "Taml": "Noto Sans Tamil",
    "Telu": "Noto Sans Telugu",
    "Knda": "Noto Sans Kannada",
    "Mlym": "Noto Sans Malayalam",
    "Bali": "Noto Sans Balinese",
    "Java": "Noto Sans Javanese",
    "Sund": "Noto Sans Sundanese",
    "Batk": "Noto Sans Batak",
    "Thai": "Noto Sans Thai",
    "Laoo": "Noto Sans Lao",
    "Khmr": "Noto Sans Khmer",
    "Mymr": "Noto Sans Myanmar",
    "Arab": "Noto Sans Arabic",
    "Hebr": "Noto Sans Hebrew",
    "Syrc": "Noto Sans Syriac",
    "Thaa": "Noto Sans Thaana",
    "Nkoo": "Noto Sans NKo",
    "Hans": "Noto Sans SC",
    "Hant": "Noto Sans TC",
    "Jpan": "Noto Sans JP",
    "Kore": "Noto Sans KR",
    "Tibt": "Noto Sans Tibetan",
    "Ethi": "Noto Sans Ethiopic",
    "Tfng": "Noto Sans Tifinagh",
    "Zsym": "Noto Color Emoji",
}


def load_test_strings() -> Dict:
    with TEST_STRINGS_PATH.open("r", encoding="utf-8") as fh:
        return json.load(fh)


def primary_font(script_value: str) -> str:
    # Handle combined scripts like "Latn+Bali" by picking the first segment.
    primary = script_value.split("+")[0]
    return FONT_MAP.get(primary, "Noto Sans")


def add_rtl_paragraph_parms(paragraph) -> None:
    p_pr = paragraph._element.get_or_add_pPr()
    bidi = OxmlElement("w:bidi")
    p_pr.append(bidi)
    paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT


def add_lang_and_font(run, font_name: str, lang: str) -> None:
    run.font.name = font_name
    r_pr = run._element.get_or_add_rPr()

    r_fonts = OxmlElement("w:rFonts")
    r_fonts.set(qn("w:ascii"), font_name)
    r_fonts.set(qn("w:hAnsi"), font_name)
    r_fonts.set(qn("w:cs"), font_name)
    r_pr.append(r_fonts)

    lang_elem = OxmlElement("w:lang")
    lang_elem.set(qn("w:val"), lang)
    r_pr.append(lang_elem)


def generate_docx(test_strings: Dict) -> Path:
    doc = Document()
    doc.add_heading("Multi-Script Reference", 0)
    doc.add_paragraph("Representative text for multiple scripts, used as fixtures.")

    for category, samples in test_strings.items():
        doc.add_heading(category.replace("_", " ").title(), level=1)
        for key, entry in samples.items():
            text = entry["text"]
            lang = entry.get("language", "und")
            script = entry.get("script", "Latn")
            desc = entry.get("description", "")
            font = primary_font(script)
            rtl = script.startswith("Arab") or script in {"Hebr", "Syrc", "Thaa", "Nkoo"} or "Arab" in script

            para = doc.add_paragraph()
            para.add_run(f"{key.replace('_', ' ').title()} ({lang}, {script})\n")
            run = para.add_run(text)
            add_lang_and_font(run, font, lang)
            if rtl:
                add_rtl_paragraph_parms(para)
            if desc:
                doc.add_paragraph(desc, style="Intense Quote")

    out_path = REFERENCE_DOCS / "multiscript.docx"
    doc.save(out_path)
    return out_path


def add_pptx_paragraph(frame, text: str, font_name: str, size: int = 22, rtl: bool = False) -> None:
    p = frame.add_paragraph()
    p.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)


def generate_pptx(test_strings: Dict) -> Path:
    if not HAS_PPTX:
        raise RuntimeError("python-pptx not installed; cannot generate PPTX")
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[5]  # blank

    # Title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(8.0), Inches(1.0))
    title_frame = title_box.text_frame
    title_frame.text = "Multi-Script Reference"
    title_frame.paragraphs[0].runs[0].font.size = Pt(36)

    for category, samples in test_strings.items():
        slide = prs.slides.add_slide(title_slide_layout)
        header = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(8.0), Inches(0.6))
        header_frame = header.text_frame
        header_frame.text = category.replace("_", " ").title()
        header_frame.paragraphs[0].runs[0].font.size = Pt(28)

        body = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.0), Inches(5.0))
        frame = body.text_frame
        frame.word_wrap = True

        for key, entry in samples.items():
            text = entry["text"]
            lang = entry.get("language", "und")
            script = entry.get("script", "Latn")
            desc = entry.get("description", "")
            font = primary_font(script)
            rtl = script.startswith("Arab") or script in {"Hebr", "Syrc", "Thaa", "Nkoo"} or "Arab" in script

            add_pptx_paragraph(frame, f"{key.replace('_', ' ').title()} ({lang}, {script})", font, size=18, rtl=rtl)
            add_pptx_paragraph(frame, text, font, size=24, rtl=rtl)
            if desc:
                add_pptx_paragraph(frame, desc, font, size=16, rtl=rtl)

    out_path = REFERENCE_DOCS / "multiscript.pptx"
    prs.save(out_path)
    return out_path


def main() -> Tuple[Path, Path]:
    reference_docs = [REFERENCE_DOCS]
    for path in reference_docs:
        path.mkdir(parents=True, exist_ok=True)

    test_strings = load_test_strings()
    docx_path = generate_docx(test_strings)
    print(f"✓ Generated {docx_path.relative_to(Path.cwd())}")

    pptx_path = None
    if HAS_PPTX:
        pptx_path = generate_pptx(test_strings)
        print(f"✓ Generated {pptx_path.relative_to(Path.cwd())}")
    else:
        print("✗ Skipped PPTX generation (python-pptx not installed)")

    return docx_path, pptx_path


if __name__ == "__main__":
    main()
