#!/usr/bin/env python3
"""
Generate minimal PPTX files for reference.

These files are used to understand the XML structure of OOXML PresentationML.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import os

OUTPUT_DIR = os.path.join(os.path.dirname(__file__),
                          "../../tests/fixtures/reference_docs")


def save(prs, name):
    path = os.path.join(OUTPUT_DIR, name)
    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    print(f"Created: {path}")


def create_minimal_pptx():
    """Single title slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    slide.shapes.title.text = "Hello World"
    slide.placeholders[1].text = "Subtitle text"
    save(prs, "minimal.pptx")


def create_content_pptx():
    """Title + bullet list content slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
    slide.shapes.title.text = "Content Slide"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "First bullet point"
    p = tf.add_paragraph()
    p.text = "Second bullet point"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Sub-bullet"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Third bullet point"
    p.level = 0
    save(prs, "content.pptx")


def create_styled_pptx():
    """Slide with styled text (bold, italic, color, size, font)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True

    # Bold
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Bold text"
    run.font.bold = True
    run.font.size = Pt(24)

    # Italic
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Italic text"
    run.font.italic = True
    run.font.size = Pt(24)

    # Colored
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Red colored text"
    run.font.color.rgb = RGBColor(255, 0, 0)
    run.font.size = Pt(24)

    # Underline
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Underlined text"
    run.font.underline = True
    run.font.size = Pt(24)

    # Different font
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Arial font"
    run.font.name = "Arial"
    run.font.size = Pt(24)

    save(prs, "styled.pptx")


def create_table_pptx():
    """Slide with a table."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank

    rows, cols = 4, 3
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(3))
    table = table_shape.table

    headers = ["Name", "Value", "Status"]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h

    data = [
        ["Alpha", "100", "Active"],
        ["Beta", "200", "Inactive"],
        ["Gamma", "300", "Active"],
    ]
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            table.cell(r + 1, c).text = val

    save(prs, "table.pptx")


def create_image_pptx():
    """Slide with an image."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank

    img_path = os.path.join(os.path.dirname(__file__),
                            "../../tests/fixtures/reference_docs/assets/inline.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(2), Inches(2), Inches(4), Inches(3))
    else:
        txBox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(4), Inches(1))
        txBox.text_frame.text = "[inline.png missing — run minimal_docx.py first]"

    save(prs, "image.pptx")


def create_chart_pptx():
    """Slides with bar, line, and pie charts."""
    prs = Presentation()

    # Bar chart
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
    chart_data.add_series("Revenue", (120, 150, 170, 200))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1.5), Inches(8), Inches(5),
        chart_data
    )

    # Line chart
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    chart_data2 = CategoryChartData()
    chart_data2.categories = ["Jan", "Feb", "Mar", "Apr", "May"]
    chart_data2.add_series("Users", (100, 150, 130, 180, 220))
    slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(1), Inches(1.5), Inches(8), Inches(5),
        chart_data2
    )

    # Pie chart
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    chart_data3 = CategoryChartData()
    chart_data3.categories = ["Product A", "Product B", "Product C"]
    chart_data3.add_series("Market Share", (45, 35, 20))
    slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(1), Inches(1.5), Inches(8), Inches(5),
        chart_data3
    )

    save(prs, "charts.pptx")


def create_notes_pptx():
    """Slide with speaker notes."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Slide with Notes"
    slide.placeholders[1].text = "Check the speaker notes."

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "These are the speaker notes for this slide."

    save(prs, "notes.pptx")


def create_two_column_pptx():
    """Two-column layout slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[3])  # Two Content layout
    slide.shapes.title.text = "Two Column Layout"
    slide.placeholders[1].text = "Left column content"
    slide.placeholders[2].text = "Right column content"
    save(prs, "two_column.pptx")


if __name__ == "__main__":
    print("Generating PPTX reference documents...")
    create_minimal_pptx()
    create_content_pptx()
    create_styled_pptx()
    create_table_pptx()
    create_image_pptx()
    create_chart_pptx()
    create_notes_pptx()
    create_two_column_pptx()
    print("Done! Generated 8 reference PPTX files.")
